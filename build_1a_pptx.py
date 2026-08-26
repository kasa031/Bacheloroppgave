#!/usr/bin/env python3
"""Build the ACIT4280 1A group presentation (20 min, all members)."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt
from pathlib import Path

FIG = Path("/workspace/figures")

NAVY = RGBColor(0x1B, 0x36, 0x5D)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFF, 0xF1, 0xD1)
LINE = RGBColor(0xD0, 0xD5, 0xDD)
GREEN = RGBColor(0x1B, 0x6B, 0x3A)
AMBER = RGBColor(0x9A, 0x5B, 0x00)
RED = RGBColor(0x8B, 0x1E, 0x1E)
PAPER = RGBColor(0xFF, 0xFE, 0xFB)

W = Inches(13.333)
H = Inches(7.5)


def set_run(run, text, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    run.font.italic = False


def add_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def add_tb(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run() if p.runs else p.runs[0] if False else _first_run(p), text, size, bold, color, font)
    return box


def _first_run(p):
    if p.runs:
        return p.runs[0]
    return p.add_run()


def textbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def p_run(tf, text, size=18, bold=False, color=INK, space_before=0, space_after=6, align=PP_ALIGN.LEFT, font="Calibri"):
    if not tf.paragraphs[0].runs and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    set_run(r, text, size, bold, color, font)
    return p


def footer(slide, n, total=18):
    add_rect(slide, 0, Inches(7.28), W, Inches(0.22), NAVY)
    add_tb(slide, Inches(0.4), Inches(7.28), Inches(10), Inches(0.22),
           "ACIT4280 Privacy by Design  |  Third-party sharing and lawful basis  |  3 Sep 2026",
           size=10, color=WHITE)
    add_tb(slide, Inches(11.4), Inches(7.28), Inches(1.5), Inches(0.22),
           f"{n} / {total}", size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header_bar(slide, kicker, title, speaker=""):
    add_rect(slide, 0, 0, W, Inches(1.15), NAVY)
    add_rect(slide, 0, Inches(1.15), W, Inches(0.08), CREAM)
    add_tb(slide, Inches(0.5), Inches(0.12), Inches(8.2), Inches(0.32),
           kicker, size=12, color=CREAM, bold=True)
    if speaker:
        add_tb(slide, Inches(8.8), Inches(0.12), Inches(4.0), Inches(0.32),
               speaker, size=12, color=CREAM, bold=True, align=PP_ALIGN.RIGHT)
    add_tb(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.62),
           title, size=28, color=WHITE, bold=True)


def bullet_slide(prs, n, kicker, title, bullets, note, sizes=None, speaker=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, kicker, title, speaker=speaker)
    tf = textbox(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.5))
    for i, b in enumerate(bullets):
        sz = (sizes[i] if sizes else 20)
        p_run(tf, b, size=sz, space_after=10)
    footer(s, n)
    notes(s, note)
    return s


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    total = 18

    # 1 title
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, Inches(5.85), W, Inches(1.65), CREAM)
    add_tb(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
           "ACIT4280 Privacy by Design",
           size=16, color=CREAM, bold=True)
    add_tb(s, Inches(0.7), Inches(2.05), Inches(12), Inches(1.8),
           "Third-party data sharing, tracking\nand lawful basis on Norwegian sites",
           size=36, color=WHITE, bold=True)
    add_tb(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.9),
           "Humna Akhtar  ·  Mithun Chandra Debnath\nKarina Sætersdal Nilssen  ·  Sumit Prasad Sah",
           size=18, color=WHITE)
    add_tb(s, Inches(0.7), Inches(6.1), Inches(12), Inches(1.0),
           "Oslo Metropolitan University  ·  20 minutes  ·  All members present\nDue 3 September 2026, 08:00",
           size=16, color=NAVY)
    notes(s, "SPEAKER: Humna (open). Welcome the audience. State the two parts: Webbkoll on 18 sites, then ICO on five privacy notices.")

    # 2 agenda
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "20 minutes  ·  four speakers", "Who presents what", speaker="Humna")
    rows = [
        ("Humna Akhtar", "Task, method, Webbkoll rules", "~5 min"),
        ("Mithun Chandra Debnath", "Webbkoll findings: highest/lowest, sectors, klassekampen", "~5 min"),
        ("Karina Sætersdal Nilssen", "ICO tool, one purpose, Skatteetaten, Netflix", "~5 min"),
        ("Sumit Prasad Sah", "fotball.no, document.no, babyshop.no, close", "~5 min"),
    ]
    y = Inches(1.55)
    for name, job, mins in rows:
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.15), CREAM)
        add_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.15), NAVY)
        add_tb(s, Inches(0.85), y + Inches(0.12), Inches(8.5), Inches(0.4), name, size=20, bold=True, color=NAVY)
        add_tb(s, Inches(0.85), y + Inches(0.52), Inches(9.5), Inches(0.5), job, size=16, color=INK)
        add_tb(s, Inches(10.5), y + Inches(0.35), Inches(2.1), Inches(0.45), mins, size=16, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        y += Inches(1.3)
    footer(s, 2, total)
    notes(s, "SPEAKER: Humna. Read the split once. Each person stays on their slides. Hand over by name.")

    # 3 two parts
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "Structure", "The report has two parts", speaker="Humna")
    add_rect(s, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.0), CREAM)
    add_tb(s, Inches(0.75), Inches(1.75), Inches(5.4), Inches(0.4), "Part 1  ·  Webbkoll", size=22, bold=True, color=NAVY)
    tf = textbox(s, Inches(0.75), Inches(2.3), Inches(5.4), Inches(4.0))
    for line in [
        "18 Norwegian-facing sites",
        "Shopping, government, news/media, sport",
        "Server location, cookies, 3rd-party requests, countries excl. Norway",
        "Graphics per sector",
        "Highest and lowest contact; most countries; restrained e-government",
    ]:
        p_run(tf, "•  " + line, size=16, space_after=8)
    add_rect(s, Inches(6.9), Inches(1.55), Inches(5.9), Inches(5.0), NAVY)
    add_tb(s, Inches(7.15), Inches(1.75), Inches(5.4), Inches(0.4), "Part 2  ·  ICO lawful basis", size=22, bold=True, color=CREAM)
    tf = textbox(s, Inches(7.15), Inches(2.3), Inches(5.4), Inches(4.0))
    for line in [
        "Five privacy notices",
        "skatteetaten, Netflix, fotball.no, document.no, babyshop.no",
        "One concrete purpose per ICO run",
        "Compare notice vs official ICO report",
        "Cookies, membership and checkout stay separate purposes",
    ]:
        p_run(tf, "•  " + line, size=16, color=WHITE, space_after=8)
    footer(s, 3, total)
    notes(s, "SPEAKER: Humna. Stress one purpose at a time. Cookie clicks are not the lawful basis for tax or checkout.")

    # 4 method
    bullet_slide(
        prs, 4, "Part 1", "How we measured with Webbkoll",
        [
            "English Webbkoll interface",
            "One live Chromium visit: no add-ons, no Do Not Track, no consent click",
            "Scan clock on the result page is part of the measurement",
            "Server country: KeyCDN Country field for the IP",
            "Norway dropped from the shared-country list",
            "If KeyCDN returns no country: recorded as not available",
            "Two tables: Table 2 first measurement, Table 3 repeat measurement",
        ],
        "SPEAKER: Humna. Webbkoll does not read a privacy notice and does not decide Article 6. Request volume is not cookie volume.",
        sizes=[18]*7,
        speaker="Humna",
    )

    # 5 rules
    bullet_slide(
        prs, 5, "Part 1", "What Webbkoll does not do",
        [
            "It does not click Ja/Nei or Cookiebot.",
            "Zero third-party cookies can still mean many third-party requests.",
            "ikea.no is the only row with external cookies (2), after redirect to ikea.com.",
            "The two measurements are not averaged. A later visit can differ (24-hour store).",
            "Article 6 is not on the Webbkoll page. That is the ICO tool.",
        ],
        "SPEAKER: Humna. Hand over to Mithun for the numbers.",
        speaker="Humna",
    )

    # 6 highest lowest
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "Part 1  ·  Table 4", "Highest and lowest third-party requests", speaker="Mithun")
    # two columns of cards
    high = [
        ("1  fotball.no", "Sport  ·  113 requests  ·  5 countries"),
        ("2  worldofwarcraft.com", "News/media  ·  90  ·  4"),
        ("3  document.no", "News/media  ·  51  ·  6 (widest list)"),
        ("4  aftenposten.no", "News/media  ·  48  ·  5"),
        ("5  klassekampen.no", "News/media  ·  46  ·  2"),
    ]
    low = [
        ("1  lanekassen.no", "Government  ·  1 request"),
        ("2  altinn.no", "Government  ·  5"),
        ("3  nrk.no", "News/media  ·  7"),
        ("4  spotify.no", "News/media  ·  10"),
        ("5  skiforeningen.no", "Sport  ·  14"),
    ]
    add_tb(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4), "Highest five", size=18, bold=True, color=NAVY)
    add_tb(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.4), "Lowest five", size=18, bold=True, color=NAVY)
    y = Inches(1.85)
    for (a, b), (c, d) in zip(high, low):
        add_rect(s, Inches(0.5), y, Inches(5.9), Inches(0.9), CREAM)
        add_tb(s, Inches(0.7), y + Inches(0.08), Inches(5.5), Inches(0.35), a, size=16, bold=True, color=NAVY)
        add_tb(s, Inches(0.7), y + Inches(0.42), Inches(5.5), Inches(0.35), b, size=14, color=INK)
        add_rect(s, Inches(7.0), y, Inches(5.8), Inches(0.9), RGBColor(0xF4, 0xF6, 0xF8))
        add_tb(s, Inches(7.2), y + Inches(0.08), Inches(5.4), Inches(0.35), c, size=16, bold=True, color=NAVY)
        add_tb(s, Inches(7.2), y + Inches(0.42), Inches(5.4), Inches(0.35), d, size=14, color=INK)
        y += Inches(1.0)
    footer(s, 6, total)
    notes(s, "SPEAKER: Mithun. Rank is third-party REQUESTS on Table 2, not cookies. The repeat measurement puts babyshop.no at 101; rankings follow the first measurement. skatteetaten.no also 14 requests; skiforeningen listed because fewer countries.")

    # 7 findings + charts
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "Part 1  ·  Figures 2 to 4", "Third-party requests by site and sector", speaker="Mithun")
    s.shapes.add_picture(str(FIG / "fig3_requests_all18.png"), Inches(0.3), Inches(1.38), Inches(6.5), Inches(5.55))
    s.shapes.add_picture(str(FIG / "fig5_share_requests_pie.png"), Inches(6.95), Inches(1.38), Inches(6.0), Inches(2.7))
    s.shapes.add_picture(str(FIG / "fig6_requests_per_sector.png"), Inches(6.95), Inches(4.15), Inches(6.0), Inches(2.75))
    footer(s, 7, total)
    notes(s, "SPEAKER: Mithun. Point to fotball.no at 113 and lanekassen.no at 1. Request volume is not cookie volume.")

    # 8 klassekampen
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "Part 1  ·  method example", "klassekampen.no: requests without third-party cookies", speaker="Mithun")
    s.shapes.add_picture(str(FIG / "fig2_webbkoll_results_klassekampen.png"), Inches(0.4), Inches(1.35), Inches(12.5), Inches(2.55))
    s.shapes.add_picture(str(FIG / "fig4_keycdn_lookup_klassekampen.png"), Inches(2.4), Inches(4.0), Inches(8.5), Inches(3.05))
    footer(s, 8, total)
    notes(s, "SPEAKER: Mithun. This shows that request volume is not cookie volume. Hand over to Karina for ICO.")

    # 9 ICO intro
    bullet_slide(
        prs, 9, "Part 2", "ICO lawful basis: one purpose at a time",
        [
            "Tool: ICO Lawful basis interactive guidance (updated 14 April 2026).",
            "Five services, one concrete processing activity per run.",
            "Find the privacy notice, run the tool, save the official Word report.",
            "Compare: what the notice claims vs what ICO marks APPROPRIATE.",
            "Do not mix cookies, marketing, membership and checkout in one run.",
            "UK tool also has recognised legitimate interest (five narrow public purposes). None of our five purposes fitted that list.",
        ],
        "SPEAKER: Karina. Hold up that we used official ICO reports dated 26 August 2026.",
        sizes=[18]*6,
        speaker="Karina",
    )

    # 10 art 6 vs cookies
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "Part 2", "Article 6 is not a cookie click", speaker="Karina")
    add_rect(s, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.0), CREAM)
    add_tb(s, Inches(0.75), Inches(1.8), Inches(5.4), Inches(0.5), "Article 6 (GDPR)", size=22, bold=True, color=NAVY)
    tf = textbox(s, Inches(0.75), Inches(2.45), Inches(5.4), Inches(3.8))
    for line in [
        "Why may we process this personal data?",
        "Contract, legal obligation, public task, consent, legitimate interests, vital interests",
        "Example: tax and folkeregister under named acts",
    ]:
        p_run(tf, "•  " + line, size=16, space_after=10)
    add_rect(s, Inches(6.9), Inches(1.55), Inches(5.9), Inches(5.0), NAVY)
    add_tb(s, Inches(7.15), Inches(1.8), Inches(5.4), Inches(0.5), "ePrivacy / cookies", size=22, bold=True, color=CREAM)
    tf = textbox(s, Inches(7.15), Inches(2.45), Inches(5.4), Inches(3.8))
    for line in [
        "May we store/read info on the device?",
        "Necessary cookies: often no extra consent",
        "Analytics and ads: consent (Ja/Nei, Cookiebot)",
        "Nei test 24 Aug 2026 on skatteetaten.no/person/",
    ]:
        p_run(tf, "•  " + line, size=16, color=WHITE, space_after=10)
    footer(s, 10, total)
    notes(s, "SPEAKER: Karina. A clean ICO run with Consent Q1 = No marks consent NOT APPROPRIATE for tax. Cookies are a separate purpose and a separate run.")

    # 11 table 5
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "Part 2  ·  Table 5", "ICO result vs the notice", speaker="Karina")
    data = [
        ("Service", "Purpose", "ICO", "Match"),
        ("skatteetaten.no", "Tax / folkeregister", "Legal obligation + public task", "Yes"),
        ("skatteetaten.no", "Optional cookies", "Consent INCONCLUSIVE; none APPROPRIATE", "Partial"),
        ("netflix.no", "Paid streaming", "Contract", "Yes"),
        ("fotball.no", "Name + club, 13+", "Legitimate interests", "Yes"),
        ("document.no", "Google Signals", "Consent INCONCLUSIVE; none APPROPRIATE", "Partial"),
        ("babyshop.no", "Checkout / delivery", "Contract", "Yes"),
    ]
    left = Inches(0.4)
    top = Inches(1.45)
    widths = [Inches(2.6), Inches(3.2), Inches(5.3), Inches(1.4)]
    row_h = Inches(0.72)
    x = left
    for w, htxt in zip(widths, data[0]):
        add_rect(s, x, top, w, row_h, NAVY)
        add_tb(s, x + Inches(0.08), top + Inches(0.18), w - Inches(0.1), Inches(0.4), htxt, size=13, bold=True, color=WHITE)
        x += w
    for i, row in enumerate(data[1:]):
        y = top + row_h * (i + 1)
        bg = CREAM if row[3] == "Partial" else (RGBColor(0xF4, 0xF6, 0xF8) if i % 2 else WHITE)
        x = left
        colors = [INK, INK, INK, GREEN if row[3] == "Yes" else AMBER]
        bolds = [True, False, False, True]
        for j, (w, cell) in enumerate(zip(widths, row)):
            add_rect(s, x, y, w, row_h, bg)
            add_tb(s, x + Inches(0.08), y + Inches(0.18), w - Inches(0.12), Inches(0.45), cell, size=13, bold=bolds[j], color=colors[j])
            x += w
    footer(s, 11, total)
    notes(s, "SPEAKER: Karina. Four of five core purposes match. Two consent tests fail: skatteetaten cookies (Q6, banner too short) and document.no Signals (request not separate from terms).")

    # 12 skatteetaten
    bullet_slide(
        prs, 12, "ICO  ·  government", "skatteetaten.no: law for tax, consent for cookies",
        [
            "Purpose 1: identity, income and tax data for tax and folkeregister.",
            "Notice: first and foremost laid down in law. Opt-out generally not possible.",
            "Named acts: skatteforvaltningsloven, folkeregisterloven, skattebetalingsloven.",
            "Controller: the Director General of Taxation (Skattedirektøren).",
            "ICO: legal obligation APPROPRIATE, public task APPROPRIATE.",
            "Consent NOT APPROPRIATE for core tax (no genuine free choice). Contract and LI NOT APPROPRIATE.",
            "Purpose 2: optional stats cookies (Skyra, Matomo, GA, Siteimprove). Ja/Nei. ICO: consent INCONCLUSIVE (Q6). No basis APPROPRIATE.",
        ],
        "SPEAKER: Karina. Two purposes, two reports. Tax: law. Cookies: they aim at consent, but ICO says name controller, purpose and types more clearly. Nei test 24 August 2026 on skatteetaten.no/person/.",
        sizes=[16]*7,
        speaker="Karina",
    )

    # 13 netflix
    bullet_slide(
        prs, 13, "ICO  ·  news / media", "netflix.no: contract for the paid service",
        [
            "Purpose: account and payment data to provide the subscription.",
            "EEA/UK notice: contractual necessity to provide the service to members.",
            "ICO: Contract APPROPRIATE. All other bases NOT APPROPRIATE.",
            "Consent: not appropriate / likely invalid (no real ongoing choice).",
            "LI, legal obligation and consent in the notice apply to other purposes.",
            "Keep ads, recommendations and marketing out of this run.",
        ],
        "SPEAKER: Karina. Netflix requires 18 or a parent. Hand over to Sumit.",
        sizes=[18]*6,
        speaker="Karina",
    )

    # 14 fotball
    bullet_slide(
        prs, 14, "ICO  ·  sport", "fotball.no: legitimate interests for match history",
        [
            "Purpose: publish name and club of active players 13+ (kamphistorikk).",
            "Notice: public-interest match history, opt-out via the club.",
            "ICO: Legitimate interests APPROPRIATE after necessity + balancing.",
            "NFF is a sports federation, not a public authority (public task No).",
            "FIKS membership can use contract. That is a different purpose.",
            "Stamdata 'samtykke om publisering' does not match ICO consent (likely invalid).",
            "Cookiebot (Deny on 26 Aug 2026) is also a different purpose.",
        ],
        "SPEAKER: Sumit. ICO also points to an LIA and children's-data guidance for 13-17.",
        sizes=[17]*7,
        speaker="Sumit",
    )

    # 15 document
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, PAPER)
    header_bar(s, "ICO  ·  news  ·  the mismatch", "document.no: Google Signals", speaker="Sumit")
    add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), CREAM)
    add_tb(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.9),
           "ICO marks no basis APPROPRIATE. Consent is INCONCLUSIVE.",
           size=22, bold=True, color=RED)
    tf = textbox(s, Inches(0.55), Inches(2.9), Inches(12.2), Inches(4.0))
    for line in [
        "Purpose: Google Signals / ad analytics (demographics, interests, cross-device).",
        "Notice does not name Article 6. Runs only if logged into Google with ad personalization.",
        "Opt-out: Google Ads Settings, mobile, NAI, Analytics add-on.",
        "Q5: consent request is not clear, prominent and separate from terms.",
        "Choice sits at Google. Pluss terms also say you accept the privacy notice by using the service.",
        "Contract and LI correctly rejected for this advertising purpose.",
        "Match: partial. They aim at consent; ICO says fix the request.",
    ]:
        p_run(tf, "•  " + line, size=17, space_after=8)
    footer(s, 15, total)
    notes(s, "SPEAKER: Sumit. This is the strongest 1A finding: a notice can talk like consent and still fail the ICO consent standard.")

    # 16 babyshop
    bullet_slide(
        prs, 16, "ICO  ·  shopping", "babyshop.no: contract for checkout",
        [
            "Purpose: name, address, contact, order and payment to deliver goods.",
            "Terms point 1: a sales contract (18 years or guardian).",
            "ICO: Contract APPROPRIATE. Other bases NOT APPROPRIATE.",
            "Notice does not label Art. 6(1)(b). Section 5 is profiling, not purchase.",
            "Marketing/profiling: consent (separate). Cookie policy: necessary vs measurement vs marketing.",
            "Bundled 'consent' to cookies inside the terms is not valid Art. 6 consent.",
        ],
        "SPEAKER: Sumit. Same pattern as Netflix: contract for the core service, consent for ads.",
        sizes=[18]*6,
        speaker="Sumit",
    )

    # 17 takeaways
    bullet_slide(
        prs, 17, "Close", "What we want the room to remember",
        [
            "Webbkoll measures contact, not lawfulness.",
            "One purpose per ICO run, or the tool mixes bases.",
            "fotball.no has the highest request count and still zero third-party cookies.",
            "Four of five core ICO purposes match the notice.",
            "Two consent tests fail: skatteetaten cookies (Q6) and document.no Signals (Q5).",
            "Legal obligation (tax) is not the same as cookie consent (Ja/Nei).",
        ],
        "SPEAKER: Sumit. Then questions. If asked about access requests: that is outside the scope of this report.",
        sizes=[20]*6,
        speaker="Sumit",
    )

    # 18 questions
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, NAVY)
    add_tb(s, Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
           "Questions", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, Inches(0.7), Inches(3.8), Inches(12), Inches(1.4),
           "Humna  ·  Mithun  ·  Karina  ·  Sumit\nACIT4280 Privacy by Design",
           size=20, color=CREAM, align=PP_ALIGN.CENTER)
    notes(s, "SPEAKER: Sumit, then any member. Likely questions: why not test Document Pluss as contract (would duplicate Netflix); why LI for football names (notice + opt-out + limited data).")

    out = "/workspace/ACIT4280_1A_presentation.pptx"
    prs.save(out)
    print("Wrote", out, "slides", len(prs.slides))


if __name__ == "__main__":
    main()
